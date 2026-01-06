import os
import re
from typing import List, Set, Tuple

import streamlit as st
import torch
from pptx import Presentation
from docx import Document
from transformers import (
    AutoTokenizer,
    AutoModelForSeq2SeqLM,
    AutoModelForQuestionAnswering,
    pipeline,
)

# ==========================
# PAGE CONFIG + STYLING
# ==========================

st.set_page_config(
    page_title="Local Q&A Studio",
    layout="wide",
)

st.markdown(
    """
    <style>
    body {
        background: radial-gradient(circle at top left, #1e293b, #020617);
    }
    .stApp {
        background: radial-gradient(circle at top left, #1e293b, #020617);
        color: #e5e7eb;
        font-family: -apple-system, BlinkMacSystemFont, "SF Pro Text",
                     system-ui, sans-serif;
    }
    h1 {
        text-align: center;
        font-weight: 700;
        letter-spacing: 0.04em;
        margin-bottom: 0.5rem;
    }
    h2, h3 {
        font-weight: 600;
        letter-spacing: 0.02em;
    }

    .stButton>button {
        background: linear-gradient(135deg, #f97316, #ec4899);
        color: white;
        border: none;
        border-radius: 999px;
        padding: 0.45rem 1.5rem;
        font-weight: 600;
        box-shadow: 0 8px 18px rgba(0,0,0,0.4);
        transition: transform 0.06s ease-out, box-shadow 0.06s ease-out;
    }
    .stButton>button:hover {
        transform: translateY(-1px);
        box-shadow: 0 14px 28px rgba(0,0,0,0.5);
    }

    [data-baseweb="slider"] > div > div {
        background-color: #f97316 !important;
    }

    .metric-card {
        padding: 0.85rem 1rem;
        border-radius: 18px;
        background: rgba(15, 23, 42, 0.95);
        border: 1px solid rgba(148, 163, 184, 0.4);
        box-shadow: 0 10px 25px rgba(15,23,42,0.7);
        font-size: 0.9rem;
    }

    section[data-testid="stSidebar"] {
        background: rgba(15, 23, 42, 0.95);
        border-right: 1px solid rgba(148,163,184,0.3);
    }
    </style>
    """,
    unsafe_allow_html=True,
)

# ==========================
# BASIC CONFIG
# ==========================

# QG backbones
QG_MODEL_T5_PATH = "t5-final"          # your t5 model folder
QG_MODEL_FLAN_PATH = "flan_t5_final"   # your flan-t5 model folder

# QA backbones
QA_MODEL_DEBERTA_PATH = "deberta_qa_final"
QA_MODEL_DISTIL_PATH = "distilbert_qa_final"

MIN_WORDS_PER_PARAGRAPH = 20           # hardcoded best guess
NUM_QUESTIONS_PER_PARAGRAPH = 4        # hardcoded best guess
MAX_QUESTIONS_PER_PARAGRAPH = 5        # safety upper bound

DEVICE_ID = 0 if torch.cuda.is_available() else -1

# ==========================
# TEXT CLEANING
# ==========================

BULLET_CHARS = [
    "•", "◦", "●", "○", "▪", "‣", "▸", "►", "■",
    "–", "—", "·", "•\u00a0"
]

BAD_Q_PREFIXES = [
    "is this", "is it", "do you", "can you",
    "what is this passage", "what does this passage",
    "what is the passage", "how do you feel",
]


def clean_text(text: str) -> str:
    if not text:
        return ""

    for ch in BULLET_CHARS:
        text = text.replace(ch, "- ")

    text = "".join(ch for ch in text if ch == "\n" or ch.isprintable())

    lines = []
    for line in text.splitlines():
        line = line.strip()
        line = re.sub(r"^(-\s*){2,}", "- ", line)
        lines.append(line)
    text = "\n".join(lines)

    text = re.sub(r"[ \t]+", " ", text)

    return text.strip()


# ==========================
# QUALITY FILTERS
# ==========================

def is_good_question(q: str, passage: str) -> bool:
    """Heuristic filter to drop low-quality questions."""
    q = q.strip()
    words = q.split()
    if len(words) < 5 or len(words) > 16:
        return False

    if not q.endswith("?"):
        return False

    q_l = q.lower()
    for bad in BAD_Q_PREFIXES:
        if q_l.startswith(bad):
            return False

    # lexical overlap with passage content words
    passage_words = set(
        w.lower().strip(".,!?;:()[]\"'")
        for w in passage.split()
        if len(w) > 3
    )
    q_words = set(
        w.lower().strip(".,!?;:()[]\"'")
        for w in words
        if len(w) > 3
    )

    if not passage_words or not q_words:
        return False

    overlap = len(passage_words & q_words) / len(q_words)
    if overlap < 0.3:
        return False

    return True


def is_good_answer(ans: str, passage: str) -> bool:
    """Heuristic filter to keep only grounded, reasonably short answers."""
    ans = ans.strip()
    if not ans:
        return False

    word_count = len(ans.split())
    if word_count > 25:
        return False

    passage_l = passage.lower()
    ans_l = ans.lower()
    if ans_l not in passage_l:
        stripped = ans.strip(".,!?;:()[]\"'").lower()
        if stripped not in passage_l:
            return False

    return True


# ==========================
# MODEL LOADING (QG)
# ==========================

@st.cache_resource(show_spinner=True)
def load_qg_pipeline_t5():
    tok = AutoTokenizer.from_pretrained(QG_MODEL_T5_PATH)
    model = AutoModelForSeq2SeqLM.from_pretrained(QG_MODEL_T5_PATH)
    pipe_qg = pipeline(
        "text2text-generation",
        model=model,
        tokenizer=tok,
        device=DEVICE_ID,
    )
    return pipe_qg


@st.cache_resource(show_spinner=True)
def load_qg_pipeline_flan():
    tok = AutoTokenizer.from_pretrained(QG_MODEL_FLAN_PATH)
    model = AutoModelForSeq2SeqLM.from_pretrained(QG_MODEL_FLAN_PATH)
    pipe_qg = pipeline(
        "text2text-generation",
        model=model,
        tokenizer=tok,
        device=DEVICE_ID,
    )
    return pipe_qg


def get_qg_pipeline(qg_backend: str):
    if qg_backend == "T5":
        return load_qg_pipeline_t5()
    return load_qg_pipeline_flan()


# ==========================
# MODEL LOADING (QA)
# ==========================

@st.cache_resource(show_spinner=True)
def load_qa_pipeline_deberta():
    tok = AutoTokenizer.from_pretrained(QA_MODEL_DEBERTA_PATH)
    model = AutoModelForQuestionAnswering.from_pretrained(QA_MODEL_DEBERTA_PATH)
    pipe_qa = pipeline(
        "question-answering",
        model=model,
        tokenizer=tok,
        device=DEVICE_ID,
    )
    return pipe_qa


@st.cache_resource(show_spinner=True)
def load_qa_pipeline_distil():
    tok = AutoTokenizer.from_pretrained(QA_MODEL_DISTIL_PATH)
    model = AutoModelForQuestionAnswering.from_pretrained(QA_MODEL_DISTIL_PATH)
    pipe_qa = pipeline(
        "question-answering",
        model=model,
        tokenizer=tok,
        device=DEVICE_ID,
    )
    return pipe_qa


def get_qa_pipeline(backend: str):
    if backend == "DistilBERT":
        return load_qa_pipeline_distil()
    return load_qa_pipeline_deberta()


# ==========================
# QUESTION GENERATION
# ==========================

def generate_questions_from_passage(passage: str, n: int, qg_backend: str) -> List[str]:
    """
    Use selected T5/FLAN-T5 to generate up to N high-quality questions.
    Over-generates with sampling and filters with is_good_question.
    """
    if not passage.strip() or n <= 0:
        return []

    pipe_qg = get_qg_pipeline(qg_backend)

    prompt = (
        f"You are a strict exam writer.\n\n"
        f"From the passage below, write exactly {n} short-answer questions.\n\n"
        f"Rules:\n"
        f"1. Every question must be answerable using ONLY the passage.\n"
        f"2. Questions must be factual, not opinion-based.\n"
        f"3. Do not repeat ideas or words.\n"
        f"4. Keep questions between 5–12 words.\n"
        f"5. Do not ask yes/no questions.\n"
        f"6. Return ONLY the questions, one per line.\n\n"
        f"Passage:\n{passage}"
    )

    questions: List[str] = []
    seen: Set[str] = set()
    max_trials = 3  # how many generations to try

    for _ in range(max_trials):
        out = pipe_qg(
            prompt,
            max_length=160,
            do_sample=True,
            top_p=0.85,          # slightly more conservative
            top_k=40,
            temperature=0.7,
            repetition_penalty=1.2,
            num_return_sequences=1,
        )
        text = out[0]["generated_text"]

        lines = [ln.strip() for ln in text.splitlines() if ln.strip()]

        for line in lines:
            q = re.sub(r"^\s*(\d+[\).\:-]\s+|Q\d+[:\-]\s+)", "", line).strip()
            if not q:
                continue

            if not is_good_question(q, passage):
                continue

            norm = q.lower().rstrip("?.!")
            if norm in seen:
                continue
            seen.add(norm)

            questions.append(q)
            if len(questions) >= n:
                return questions

    return questions


def answer_question(question: str, context: str, qa_backend: str) -> str:
    if not question.strip() or not context.strip():
        return ""
    pipe_qa = get_qa_pipeline(qa_backend)
    res = pipe_qa(question=question, context=context)
    return res.get("answer", "").strip()


# ==========================
# FILE HELPERS (PPTX / DOCX / TXT)
# ==========================

def extract_text_from_pptx(path: str) -> str:
    prs = Presentation(path)
    out = []
    for slide in prs.slides:
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                txt = shape.text.strip()
                if txt:
                    parts.append(txt)
        if parts:
            out.append("\n".join(parts))
    return clean_text("\n\n".join(out).strip())


def extract_text_from_docx(path: str) -> str:
    doc = Document(path)
    paras = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    return clean_text("\n\n".join(paras).strip())


def extract_text_from_txt(path: str) -> str:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        text = f.read().strip()
    return clean_text(text)


def extract_text_any(path: str) -> str:
    ext = os.path.splitext(path)[1].lower()
    if ext == ".pptx":
        return extract_text_from_pptx(path)
    elif ext == ".docx":
        return extract_text_from_docx(path)
    elif ext == ".txt":
        return extract_text_from_txt(path)
    else:
        raise ValueError("Only .pptx, .docx, and .txt are supported here.")


def split_into_paragraphs(text: str, min_words: int) -> List[str]:
    raw_paras = [p.strip() for p in re.split(r"\n\s*\n", text) if p.strip()]
    paras = []
    for p in raw_paras:
        if len(p.split()) >= min_words:
            paras.append(p)
    return paras


# ==========================
# STATE FOR Q&A GENERATION
# ==========================

def ensure_qagen_state():
    if "paragraphs" not in st.session_state:
        st.session_state["paragraphs"] = []
    if "para_idx" not in st.session_state:
        st.session_state["para_idx"] = 0
    if "current_pairs" not in st.session_state:
        st.session_state["current_pairs"] = []
    if "qa_idx" not in st.session_state:
        st.session_state["qa_idx"] = 0
    if "qa_history" not in st.session_state:
        st.session_state["qa_history"] = []


def reset_qagen_state():
    for key in ["paragraphs", "para_idx", "current_pairs", "qa_idx", "qa_history"]:
        if key in st.session_state:
            del st.session_state[key]


def get_next_qa(qg_backend: str, qa_backend: str) -> Tuple[str, str]:
    """
    Generate the next Q/A pair on demand with quality filtering.
    Skips paragraphs that fail to produce any good Q/A.
    """
    ensure_qagen_state()
    paragraphs = st.session_state["paragraphs"]

    while True:
        # still have pairs from current paragraph
        if st.session_state["qa_idx"] < len(st.session_state["current_pairs"]):
            q, a = st.session_state["current_pairs"][st.session_state["qa_idx"]]
            st.session_state["qa_idx"] += 1
            return q, a

        # move to next paragraph
        if st.session_state["para_idx"] >= len(paragraphs):
            return "", ""

        passage = paragraphs[st.session_state["para_idx"]]
        st.session_state["para_idx"] += 1

        questions = generate_questions_from_passage(
            passage, NUM_QUESTIONS_PER_PARAGRAPH, qg_backend
        )
        pairs: List[Tuple[str, str]] = []

        for q in questions:
            ans = answer_question(q, passage, qa_backend)
            if not is_good_answer(ans, passage):
                continue
            pairs.append((q, ans))

        st.session_state["current_pairs"] = pairs
        st.session_state["qa_idx"] = 0

        if not pairs:
            # no good Q/A for this paragraph, try next one
            continue


# ==========================
# Q&A GENERATION UI
# ==========================

def qa_generation_ui(qg_backend: str, qa_backend: str):
    st.subheader("Question & Answer Generation")

    uploaded_file = st.file_uploader(
        "Upload PPTX, DOCX, or TXT",
        type=["pptx", "docx", "txt"],
        key="qagen_uploader",
    )

    col_start, col_reset = st.columns([1, 1])
    started = False

    if uploaded_file is not None and col_start.button("Start reading document"):
        tmp_dir = "tmp_uploads"
        os.makedirs(tmp_dir, exist_ok=True)
        upload_path = os.path.join(tmp_dir, uploaded_file.name)
        with open(upload_path, "wb") as f:
            f.write(uploaded_file.read())

        with st.spinner("Extracting text and splitting into paragraphs…"):
            text = extract_text_any(upload_path)
            paragraphs = split_into_paragraphs(text, MIN_WORDS_PER_PARAGRAPH)

        reset_qagen_state()
        ensure_qagen_state()
        st.session_state["paragraphs"] = paragraphs
        started = True

    if col_reset.button("Clear context & history"):
        reset_qagen_state()

    ensure_qagen_state()

    if st.session_state["paragraphs"]:
        st.markdown(
            f"*Paragraphs loaded: {len(st.session_state['paragraphs'])} "
            f"(min {MIN_WORDS_PER_PARAGRAPH} words per paragraph, "
            f"{NUM_QUESTIONS_PER_PARAGRAPH} questions per paragraph internally).*"
        )
    else:
        if not started:
            st.info("Upload a document and click **Start reading document**.")
        return

    if "qa_history" not in st.session_state:
        st.session_state["qa_history"] = []

    if st.button("Next question"):
        with st.spinner("Generating next question and answer…"):
            q, a = get_next_qa(qg_backend, qa_backend)
        if q:
            st.session_state["qa_history"].append((q, a))
        else:
            st.warning("No more high-quality questions could be generated from this document.")

    if not st.session_state["qa_history"]:
        st.info("Click **Next question** to start generating Q&A pairs.")
        return

    st.markdown("### Generated pairs")

    for idx, (q, a) in enumerate(st.session_state["qa_history"], start=1):
        st.markdown(f"**Pair {idx}**")
        col_q, col_a = st.columns(2)
        with col_q:
            st.markdown(
                '<div class="metric-card"><b>Question</b><br>' + q + "</div>",
                unsafe_allow_html=True,
            )
        with col_a:
            st.markdown(
                '<div class="metric-card"><b>Answer</b><br>' + (a if a else "(no answer)") + "</div>",
                unsafe_allow_html=True,
            )
        st.markdown("")


# ==========================
# MAIN APP
# ==========================

def main():
    st.title("Local Q&A Studio")

    qg_backend = st.sidebar.radio(
        "Question generator",
        ("T5", "FLAN-T5"),
    )

    qa_backend = st.sidebar.radio(
        "QA backbone",
        ("DeBERTa", "DistilBERT"),
    )

    with st.expander("Model paths / settings"):
        st.write(f"T5 QG model path: `{QG_MODEL_T5_PATH}`")
        st.write(f"FLAN-T5 QG model path: `{QG_MODEL_FLAN_PATH}`")
        st.write(f"DeBERTa QA model: `{QA_MODEL_DEBERTA_PATH}`")
        st.write(f"DistilBERT QA model: `{QA_MODEL_DISTIL_PATH}`")
        st.write(f"Using QG backbone: **{qg_backend}**")
        st.write(f"Using QA backbone: **{qa_backend}**")
        st.write(f"Device: {'GPU' if DEVICE_ID == 0 else 'CPU'}")
        st.write(
            f"Paragraph filter: at least {MIN_WORDS_PER_PARAGRAPH} words; "
            f"internally generates up to {NUM_QUESTIONS_PER_PARAGRAPH} questions "
            f"per paragraph and keeps only high-quality Q&A pairs."
        )

    qa_generation_ui(qg_backend, qa_backend)


if __name__ == "__main__":
    main()
